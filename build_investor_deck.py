#!/usr/bin/env python3
"""
Investor Pitch Deck Generator
Generates a professional investor pitch deck for multifamily real estate funds
using python-pptx. Configurable branding, data-driven slides, institutional formatting.

Usage:
    python3 build_investor_deck.py

Requires:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === BRAND CONSTANTS (customize these for your fund) ===
ACCENT = RGBColor(0x8A, 0x9A, 0x7B)       # #8A9A7B - Primary accent
SECONDARY = RGBColor(0xA8, 0xB8, 0x9A)    # Secondary accent
DARK_SLATE = RGBColor(0x34, 0x3D, 0x46)   # #343D46 - Primary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM_GRAY = RGBColor(0xDC, 0xD5, 0xCF)    # #DCD5CF - Backgrounds
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xF0)     # Light warm background
FONT = "Inter"

# Slide dimensions - 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Asset paths (place your logo files in ./assets/)
ASSETS = "./assets"
LOGO_DARK = os.path.join(ASSETS, "logo_dark.png")
LOGO_WHITE = os.path.join(ASSETS, "logo_white.png")
ICON_ACCENT = os.path.join(ASSETS, "icon_accent.png")
ICON_CIRCLE = os.path.join(ASSETS, "icon_circle.png")

OUTPUT = "./output/Investor_Pitch_Deck.pptx"


def add_accent_bar(slide, top=0):
    """Add the signature thin accent bar at top of slide"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_logo_watermark(slide):
    """Add small icon watermark to bottom-right"""
    if os.path.exists(ICON_ACCENT):
        slide.shapes.add_picture(
            ICON_ACCENT,
            SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.7),
            Inches(0.8), Inches(0.5)
        )


def add_section_label(slide, text, left, top):
    """Add small caps section label in accent color (e.g., 'THE OPPORTUNITY')"""
    txBox = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT


def add_title(slide, text, left, top, width=Inches(10), size=Pt(36), color=DARK_SLATE):
    """Add main slide title"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT
    return txBox


def add_body(slide, text, left, top, width=Inches(10), height=Inches(3), size=Pt(14), color=DARK_SLATE):
    """Add body text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.name = FONT
    p.line_spacing = Pt(22)
    return tf


def add_bullet_list(slide, items, left, top, width=Inches(10), height=Inches(4), size=Pt(14), color=DARK_SLATE):
    """Add a bulleted list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = FONT
        p.line_spacing = Pt(24)
        p.space_after = Pt(6)
        # Bullet character
        p.text = "\u2022  " + item
    return tf


def add_stat_box(slide, number, label, left, top, width=Inches(2.5), num_color=ACCENT):
    """Add a big stat with label below"""
    # Number
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = num_color
    p.font.name = FONT
    p.alignment = PP_ALIGN.LEFT

    # Label
    txBox2 = slide.shapes.add_textbox(left, top + Inches(0.75), width, Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.LEFT


def add_kv_row(tf, key, value, key_color=ACCENT, val_color=DARK_SLATE):
    """Add a key-value row to existing text frame"""
    p = tf.add_paragraph()
    run_k = p.add_run()
    run_k.text = key + ": "
    run_k.font.size = Pt(13)
    run_k.font.bold = True
    run_k.font.color.rgb = key_color
    run_k.font.name = FONT

    run_v = p.add_run()
    run_v.text = value
    run_v.font.size = Pt(13)
    run_v.font.color.rgb = val_color
    run_v.font.name = FONT
    p.line_spacing = Pt(22)
    p.space_after = Pt(4)


# === BUILD PRESENTATION ===
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # Blank


# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Dark background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_SLATE
bg.line.fill.background()

# Accent bar at top
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Logo
if os.path.exists(LOGO_WHITE):
    slide.shapes.add_picture(LOGO_WHITE, Inches(0.8), Inches(0.6), Inches(3.5))

# Fund name
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "MERIDIAN RESIDENTIAL FUND"
run.font.size = Pt(72)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = FONT

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "U.S. Multifamily Real Estate Investment Fund"
run2.font.size = Pt(28)
run2.font.color.rgb = ACCENT
run2.font.name = FONT

# Tagline
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Building a diversified portfolio of workforce housing assets\nacross high-growth Sun Belt markets"
p.font.size = Pt(16)
p.font.color.rgb = WARM_GRAY
p.font.name = FONT
p.line_spacing = Pt(26)

# Stats bar at bottom
stats = [
    ("$200M", "Target GAV"),
    ("13-16%", "Net IRR Target"),
    ("5-6%", "Cash Yield"),
    ("15+", "Years Experience"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + Inches(3) * i
    y = Inches(5.8)

    # Number
    txBox = slide.shapes.add_textbox(x, y, Inches(2.5), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT

    # Label
    txBox2 = slide.shapes.add_textbox(x, y + Inches(0.6), Inches(2.5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = WARM_GRAY
    p2.font.name = FONT

# Confidential footer
txBox = slide.shapes.add_textbox(Inches(0.8), SLIDE_H - Inches(0.5), Inches(8), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CONFIDENTIAL  |  For Qualified Investors Only"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0x6B, 0x72, 0x78)
p.font.name = FONT


# ============================================================
# SLIDE 2: THE OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "THE OPPORTUNITY", Inches(0.8), Inches(0.5))
add_title(slide, "A Structural Housing Shortage Is Creating\na Generational Opportunity in\nSun Belt Multifamily", Inches(0.8), Inches(0.9), width=Inches(11))

items = [
    "National housing deficit: The U.S. is short 4 to 7 million housing units, with the gap widening every year as new construction lags household formation.",
    "Sun Belt migration wave: Population growth in Sun Belt metros is 2 to 3x the national average, driven by job creation, affordability, and remote work flexibility.",
    "Homeownership out of reach: Rising mortgage rates and home prices have pushed monthly ownership costs 40%+ above renting, keeping millions in the renter pool.",
    "Class B value-add spread: Unrenovated workforce housing trades at a 150 to 250 bps cap rate spread to Class A, creating margin for renovation-driven returns.",
    "Millennial and Gen Z demand: 72 million millennials and 68 million Gen Z adults are entering peak renting years, sustaining long-term occupancy tailwinds.",
]

add_bullet_list(slide, items, Inches(0.8), Inches(2.8), width=Inches(11.5), size=Pt(14))


# ============================================================
# SLIDE 3: ABOUT THE COMPANY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "ABOUT US", Inches(0.8), Inches(0.5))
add_title(slide, "15+ Years of Multifamily Investment\nAcross the U.S. Sun Belt", Inches(0.8), Inches(0.9), width=Inches(10))

# Left column - about text
tf = add_body(slide, "", Inches(0.8), Inches(2.4), width=Inches(5.5), height=Inches(4))
p = tf.paragraphs[0]
p.text = "Meridian Capital Partners is a multifamily-focused investment firm specializing in the acquisition, renovation, and stabilization of workforce housing across high-growth Sun Belt markets."
p.font.size = Pt(15)
p.font.color.rgb = DARK_SLATE
p.font.name = FONT
p.line_spacing = Pt(24)

items_about = [
    "Full-service acquisition, renovation, and stabilization platform",
    "3,200+ units acquired and renovated across target markets",
    "In-house property management for operational control",
    "Established broker relationships in every target MSA",
    "In-house construction management capabilities",
    "Proprietary deal flow from off-market and direct seller channels",
]
for item in items_about:
    p = tf.add_paragraph()
    p.text = "\u2022  " + item
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_SLATE
    p.font.name = FONT
    p.line_spacing = Pt(22)
    p.space_after = Pt(4)

# Right column - key stats
stats_right = [
    ("15+", "Years Multifamily\nExperience"),
    ("3,200+", "Units Acquired\n& Renovated"),
    ("97%", "Average Historical\nOccupancy"),
    ("$350M+", "Total Transaction\nVolume"),
]
for i, (num, label) in enumerate(stats_right):
    x = Inches(7.5)
    y = Inches(2.4) + Inches(1.2) * i
    add_stat_box(slide, num, label, x, y, width=Inches(4.5))

# Tagline
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "\"Workforce housing. Disciplined execution. Consistent returns.\""
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = ACCENT
p.font.name = FONT


# ============================================================
# SLIDE 4: MARKET OVERVIEW - THREE COLUMNS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "MARKET OVERVIEW", Inches(0.8), Inches(0.5))
add_title(slide, "Five Target Markets Across\nthe U.S. Sun Belt", Inches(0.8), Inches(0.9), width=Inches(11))

cities = [
    {
        "name": "Austin / San Antonio",
        "color": ACCENT,
        "strengths": [
            "Fastest-growing large metro in the U.S.",
            "Major tech and defense employer base",
            "Strong rent growth with limited new Class B supply",
            "Pro-business regulatory environment",
            "Young, educated workforce driving demand",
            "Affordable relative to coastal markets",
        ]
    },
    {
        "name": "Nashville / Raleigh",
        "color": SECONDARY,
        "strengths": [
            "Top-tier job growth across healthcare and tech",
            "Population influx from Northeast and Midwest",
            "Diversified economy with Fortune 500 headquarters",
            "Strong university pipeline fueling rental demand",
            "Limited rent control or regulatory headwinds",
            "Consistent 3 to 5% annual rent growth trajectory",
        ]
    },
    {
        "name": "Tampa / Phoenix",
        "color": ACCENT,
        "strengths": [
            "Among the highest net domestic migration in the U.S.",
            "Expanding financial services and logistics sectors",
            "Deep pool of workforce housing inventory from 1980s and 1990s",
            "Pro-landlord legal framework",
            "Year-round climate attracting retirees and remote workers",
            "Strong institutional investor demand for stabilized assets",
        ]
    },
]

for i, city in enumerate(cities):
    x = Inches(0.8) + Inches(4) * i
    y_start = Inches(2.6)

    # City name
    txBox = slide.shapes.add_textbox(x, y_start, Inches(3.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = city["name"]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = city["color"]
    p.font.name = FONT

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_start + Inches(0.55), Inches(3.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = city["color"]
    line.line.fill.background()

    # Strengths
    txBox2 = slide.shapes.add_textbox(x, y_start + Inches(0.7), Inches(3.5), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for j, s in enumerate(city["strengths"]):
        if j == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "\u2022  " + s
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SLATE
        p.font.name = FONT
        p.line_spacing = Pt(20)
        p.space_after = Pt(4)


# ============================================================
# SLIDE 5: DEMAND DRIVERS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "DEMAND DRIVERS", Inches(0.8), Inches(0.5))
add_title(slide, "Why Multifamily Demand Is Accelerating\nin the Sun Belt", Inches(0.8), Inches(0.9), width=Inches(11))

drivers = [
    ("Remote Work Migration", "Remote and hybrid work policies are enabling millions of workers to relocate from high-cost coastal cities to affordable Sun Belt metros."),
    ("Housing Shortage", "The U.S. is 4 to 7 million units short. New construction remains below household formation rates, tightening supply in growth markets."),
    ("Affordability Gap", "Monthly mortgage payments now exceed rent by 40%+ in most Sun Belt metros, keeping would-be buyers in the rental market longer."),
    ("Millennial Demographics", "72 million millennials are in peak renting years. Delayed homeownership, student debt, and lifestyle preferences sustain rental demand."),
    ("Job Growth", "Sun Belt metros are adding jobs at 2 to 3x the national rate, led by healthcare, tech, logistics, and financial services employers."),
    ("Population Growth", "Sun Belt states captured over 80% of U.S. population growth since 2020, with net domestic migration accelerating post-pandemic."),
]

for i, (title, desc) in enumerate(drivers):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4) * col
    y = Inches(2.6) + Inches(2.2) * row

    # Accent box
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), Inches(0.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(x + Inches(0.2), y, Inches(3.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_SLATE
    p.font.name = FONT

    # Description
    txBox2 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.45), Inches(3.5), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.line_spacing = Pt(18)


# ============================================================
# SLIDE 6: TENANT DEMOGRAPHICS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "TENANT DEMOGRAPHICS", Inches(0.8), Inches(0.5))
add_title(slide, "Who Rents Workforce Housing\nin the Sun Belt", Inches(0.8), Inches(0.9), width=Inches(11))

segments = [
    ("Young Professionals", "Ages 25 to 34, early career", "Largest renter cohort by volume", "Relocating for job opportunities"),
    ("Healthcare Workers", "Nurses, technicians, support staff", "Stable employment, shift schedules", "Proximity to hospitals and clinics"),
    ("Tech Relocators", "Remote workers from coastal markets", "Higher income, longer lease terms", "Seeking affordability and lifestyle"),
    ("Service Industry", "Hospitality, retail, food service", "Essential workforce, consistent demand", "Price-sensitive, high turnover"),
    ("Military / Gov", "Active duty, veterans, civil servants", "Reliable income, BAH-supported rents", "Concentrated near bases and agencies"),
    ("Empty Nesters", "Ages 55+, downsizing from ownership", "Strong credit, long tenure", "Seeking low-maintenance living"),
    ("University Grads", "Recent graduates entering workforce", "Growing income trajectory", "First-time renters in metro areas"),
]

for i, (name, profile, detail, note) in enumerate(segments):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6) * col
    y = Inches(2.6) + Inches(1.1) * row

    # Segment name in accent color
    txBox = slide.shapes.add_textbox(x, y, Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT

    # Details
    txBox2 = slide.shapes.add_textbox(x + Inches(2.6), y, Inches(3), Inches(0.9))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"{profile} | {detail}"
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT

    p3 = tf2.add_paragraph()
    p3.text = note
    p3.font.size = Pt(10)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(0x6B, 0x72, 0x78)
    p3.font.name = FONT

# Bottom note
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Demographic profiles shown for demonstration purposes."
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = DARK_SLATE
p.font.name = FONT


# ============================================================
# SLIDE 7: INVESTMENT STRATEGY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "INVESTMENT STRATEGY", Inches(0.8), Inches(0.5))
add_title(slide, "Value-Add Class B Multifamily\nWith Renovation Upside", Inches(0.8), Inches(0.9), width=Inches(11))

# Left column: Acquisition Criteria
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Acquisition Criteria"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = FONT

acq_items = [
    "Class B garden-style and mid-rise apartments",
    "80 to 250 units per property",
    "1980s to 2000s vintage construction",
    "Below-market rents with renovation upside",
    "Target basis: $80K to $150K per unit",
]
add_bullet_list(slide, acq_items, Inches(0.8), Inches(3.2), width=Inches(5.5), height=Inches(3), size=Pt(13))

# Right column: Renovation Playbook
txBox = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Renovation Playbook"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = SECONDARY
p.font.name = FONT

dev_items = [
    "Interior upgrades: countertops, fixtures, flooring, appliances",
    "Exterior improvements: landscaping, paint, amenity spaces",
    "Budget: $8K to $15K per unit, scaled to market rents",
    "12 to 18 month renovation cycle per property",
    "Target rent premiums: $150 to $300 per unit post-renovation",
    "Unit turns executed on natural lease expirations",
]
add_bullet_list(slide, dev_items, Inches(7), Inches(3.2), width=Inches(5.5), height=Inches(3), size=Pt(13))

# Vertical divider
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.5), Inches(0.03), Inches(4))
line.fill.solid()
line.fill.fore_color.rgb = WARM_GRAY
line.line.fill.background()

# Portfolio target
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Target: "
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = ACCENT
r.font.name = FONT
r2 = p.add_run()
r2.text = "$200M gross asset value across 15 to 20 properties  |  80 to 250 units per deal  |  Sun Belt markets only"
r2.font.size = Pt(14)
r2.font.color.rgb = DARK_SLATE
r2.font.name = FONT


# ============================================================
# SLIDE 8: RETURN ENGINE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "RETURNS", Inches(0.8), Inches(0.5))
add_title(slide, "Five Levers Driving Returns", Inches(0.8), Inches(0.9), width=Inches(10))

levers = [
    ("1", "Below-Market Acquisition Basis", "Acquire Class B properties at a discount to replacement cost, creating embedded equity from day one."),
    ("2", "Renovation-Driven Rent Growth", "Interior and exterior upgrades generate $150 to $300 per unit rent premiums on a $8K to $15K per unit investment."),
    ("3", "Operational Efficiency", "In-house property management reduces expenses 5 to 10% versus third-party operators, improving NOI margins."),
    ("4", "Organic Market Appreciation", "Sun Belt population and job growth drive 3 to 5% annual rent increases and natural asset value appreciation."),
    ("5", "Disciplined Exit Timing", "Sell stabilized, renovated assets into strong institutional buyer demand at compressed cap rates."),
]

for i, (num, title, desc) in enumerate(levers):
    y = Inches(2.4) + Inches(0.95) * i

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.word_wrap = False
    p_c = tf_c.paragraphs[0]
    p_c.text = num
    p_c.font.size = Pt(18)
    p_c.font.bold = True
    p_c.font.color.rgb = WHITE
    p_c.font.name = FONT
    p_c.alignment = PP_ALIGN.CENTER
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), y, Inches(3.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_SLATE
    p.font.name = FONT

    # Description
    txBox2 = slide.shapes.add_textbox(Inches(5), y, Inches(7.5), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.line_spacing = Pt(18)


# ============================================================
# SLIDE 9: FUND STRUCTURE & TERMS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "FUND STRUCTURE", Inches(0.8), Inches(0.5))
add_title(slide, "Meridian Residential Fund | Key Terms", Inches(0.8), Inches(0.9), width=Inches(10))

# Two columns of terms
left_terms = [
    ("Vehicle", "Delaware LLC"),
    ("Fund Life", "7-year closed-end with two 1-year extensions"),
    ("Target GAV", "$200M across 15 to 20 properties"),
    ("Target Leverage", "55 to 65% LTV, Agency debt preferred"),
    ("Debt Profile", "Fannie Mae, Freddie Mac, and HUD financing"),
    ("Currency", "USD"),
    ("GP Commitment", "5% of fund equity"),
]

right_terms = [
    ("Target Net IRR", "13 to 16%"),
    ("Cash Yield", "5 to 6% quarterly"),
    ("Preferred Return", "7% cumulative, non-compounding"),
    ("Promote", "10% over 7% IRR\n15% over 12% IRR\n20% over 16% IRR"),
    ("AM Fee", "1.5% during investment period, then 1.0%"),
    ("Acquisition Fee", "1.0%"),
    ("Construction Mgmt", "5% of renovation budget"),
]

# Left column
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (key, val) in enumerate(left_terms):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    r_k = p.add_run()
    r_k.text = key
    r_k.font.size = Pt(13)
    r_k.font.bold = True
    r_k.font.color.rgb = ACCENT
    r_k.font.name = FONT

    p2 = tf.add_paragraph()
    p2.text = val
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.space_after = Pt(8)
    p2.line_spacing = Pt(18)

# Right column
txBox = slide.shapes.add_textbox(Inches(7), Inches(2.4), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (key, val) in enumerate(right_terms):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    r_k = p.add_run()
    r_k.text = key
    r_k.font.size = Pt(13)
    r_k.font.bold = True
    r_k.font.color.rgb = ACCENT
    r_k.font.name = FONT

    p2 = tf.add_paragraph()
    p2.text = val
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.space_after = Pt(8)
    p2.line_spacing = Pt(18)

# Vertical divider
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.4), Inches(0.03), Inches(4.5))
line.fill.solid()
line.fill.fore_color.rgb = WARM_GRAY
line.line.fill.background()


# ============================================================
# SLIDE 10: RISK MANAGEMENT
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "RISK MANAGEMENT", Inches(0.8), Inches(0.5))
add_title(slide, "Disciplined Risk Framework", Inches(0.8), Inches(0.9), width=Inches(10))

risks = [
    ("Interest Rate", "Fixed-rate Agency debt with 5 to 10 year terms; rate caps on any floating-rate exposure; DSCR floors at 1.25x"),
    ("Geographic Concentration", "No single MSA exceeds 30% of portfolio GAV; diversified across five Sun Belt markets with distinct economic drivers"),
    ("Tenant Credit", "Workforce housing targets employed renters with stable incomes; 97% historical occupancy demonstrates consistent demand"),
    ("Renovation Execution", "In-house construction management; GMP contracts with licensed contractors; contingency reserves on every project"),
    ("Insurance / Casualty", "Comprehensive property and liability coverage; flood and wind policies in applicable markets; replacement cost endorsements"),
    ("Regulatory / Rent Control", "Target markets selected for pro-landlord legal frameworks; no current rent control legislation in any target MSA"),
]

for i, (risk, mitigant) in enumerate(risks):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6) * col
    y = Inches(2.4) + Inches(1.5) * row

    # Risk category - accent
    txBox = slide.shapes.add_textbox(x, y, Inches(5.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = risk
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT

    # Mitigant
    txBox2 = slide.shapes.add_textbox(x, y + Inches(0.4), Inches(5.3), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = mitigant
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.line_spacing = Pt(18)


# ============================================================
# SLIDE 11: COMPETITIVE ADVANTAGE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_accent_bar(slide)
add_logo_watermark(slide)

add_section_label(slide, "COMPETITIVE ADVANTAGE", Inches(0.8), Inches(0.5))
add_title(slide, "Why Meridian Residential Fund", Inches(0.8), Inches(0.9), width=Inches(10))

advantages = [
    ("Local Market Knowledge", "On-the-ground teams in every target MSA with deep broker relationships, submarket expertise, and real-time deal flow intelligence."),
    ("In-House Management", "Vertically integrated property management reduces operating costs 5 to 10% versus third-party operators and improves tenant retention."),
    ("Renovation Expertise", "Proven playbook across 3,200+ units with standardized scopes, vetted contractor networks, and predictable cost and timeline execution."),
    ("Data-Driven Underwriting", "Proprietary models scoring rent comps, renovation ROI, submarket growth, and tenant demand across every target market."),
    ("Aligned Incentives", "GP commits 5% alongside LPs. Transparent fee structure. Conservative leverage. Institutional reporting and governance standards."),
]

for i, (title, desc) in enumerate(advantages):
    y = Inches(2.2) + Inches(1.0) * i

    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y + Inches(0.05), Inches(0.08), Inches(0.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.1), y, Inches(3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_SLATE
    p.font.name = FONT

    # Description
    txBox2 = slide.shapes.add_textbox(Inches(4.5), y, Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_SLATE
    p2.font.name = FONT
    p2.line_spacing = Pt(18)


# ============================================================
# SLIDE 12: THE ASK
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Dark background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_SLATE
bg.line.fill.background()

# Accent bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_section_label(slide, "THE ASK", Inches(0.8), Inches(0.5))
# Override color for dark bg
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        for p in shape.text_frame.paragraphs:
            if p.text == "THE ASK":
                p.font.color.rgb = ACCENT

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Invest in America's\nWorkforce Housing Shortage"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p.line_spacing = Pt(52)

# Key terms summary
terms_summary = [
    ("Target Net IRR", "13 to 16% (net to LPs)"),
    ("Cash Yield", "5 to 6% quarterly distributions"),
    ("Fund Life", "7-year closed-end"),
    ("Preferred Return", "7% cumulative"),
    ("GP Commitment", "5%, fully aligned"),
]

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (key, val) in enumerate(terms_summary):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    r1 = p.add_run()
    r1.text = key + "  "
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r1.font.name = FONT

    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(14)
    r2.font.color.rgb = WHITE
    r2.font.name = FONT
    p.line_spacing = Pt(28)

# Governance highlights
txBox = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(5.5), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Governance & Reporting"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = FONT

gov_items = [
    "GP-controlled IC with independent member",
    "LP advisory committee for major decisions",
    "Quarterly institutional reporting + annual GAAP audit",
    "Quarterly NAV updates",
    "Disposition-ready data rooms maintained",
]
for item in gov_items:
    p = tf.add_paragraph()
    p.text = "\u2022  " + item
    p.font.size = Pt(12)
    p.font.color.rgb = WARM_GRAY
    p.font.name = FONT
    p.line_spacing = Pt(20)

# Confidential
txBox = slide.shapes.add_textbox(Inches(0.8), SLIDE_H - Inches(0.5), Inches(8), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CONFIDENTIAL  |  For Qualified Investors Only"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0x6B, 0x72, 0x78)
p.font.name = FONT


# ============================================================
# SLIDE 13: CONTACT
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Dark background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_SLATE
bg.line.fill.background()

# Accent bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Logo
if os.path.exists(LOGO_WHITE):
    slide.shapes.add_picture(LOGO_WHITE, Inches(4.5), Inches(1.5), Inches(4.5))

# Contact info
txBox = slide.shapes.add_textbox(Inches(3), Inches(3.2), Inches(7), Inches(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "James Porter"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "CEO, Meridian Capital Partners"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT
p2.font.name = FONT
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(20)

p3 = tf.add_paragraph()
p3.text = "jporter@meridiancapitalpartners.com"
p3.font.size = Pt(14)
p3.font.color.rgb = WARM_GRAY
p3.font.name = FONT
p3.alignment = PP_ALIGN.CENTER

p4 = tf.add_paragraph()
p4.text = "meridiancapitalpartners.com"
p4.font.size = Pt(14)
p4.font.color.rgb = WARM_GRAY
p4.font.name = FONT
p4.alignment = PP_ALIGN.CENTER

p5 = tf.add_paragraph()
p5.text = "Austin, TX"
p5.font.size = Pt(14)
p5.font.color.rgb = WARM_GRAY
p5.font.name = FONT
p5.alignment = PP_ALIGN.CENTER

# Tagline
txBox = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "\"Workforce housing. Disciplined execution. Consistent returns.\""
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = ACCENT
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

# Confidential
txBox = slide.shapes.add_textbox(Inches(0.8), SLIDE_H - Inches(0.5), Inches(8), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CONFIDENTIAL  |  For Qualified Investors Only"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0x6B, 0x72, 0x78)
p.font.name = FONT


# === SAVE ===
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Deck saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
